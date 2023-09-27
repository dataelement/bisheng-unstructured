import json

import pptx
from lxml import etree


def _order_shapes(shapes):
    """Orders the shapes from top to bottom and left to right."""
    return sorted(shapes, key=lambda x: (x.top or 0, x.left or 0))


def norm_color(color):
    if not color:
        return "000000"

    if isinstance(color, list):
        return "".join(["%02X" % (i) for i in color])
    elif isinstance(color, str):
        return color
    elif isinstance(color, pptx.dml.color.RGBColor):
        return str(color)


def norm_unit(b):
    if not b:
        return b

    unit = 12700.0
    if isinstance(b, int):
        return b / unit

    return [e / unit for e in b]


def norm_align(a):
    if not a:
        return a

    n = int(a)
    type_map = {1: "left", 2: "center", 3: "right", 4: "just", 5: "dist"}
    return type_map.get(n, None)


def group_runs(runs, merge_empty=False):
    new_runs = []
    current = runs[0]
    for r in runs[1:]:
        if (
            current["font_name"] == r["font_name"]
            and current["font_size"] == r["font_size"]
            and current["font_color"] == r["font_color"]
        ):
            current["text"] += r["text"]
        else:
            if current:
                new_runs.append(current)
                current = r
            else:
                current = r
    if current:
        new_runs.append(current)

    is_empty = lambda x: len(x.strip()) == 0
    if merge_empty:
        runs = new_runs
        news_runs2 = []
        current = runs[0]
        is_c_empty = is_empty(current["text"])
        for r in runs[1:]:
            # ignore font name and size color for blank character
            # print('runs', r, is_c_empty, runs)
            if is_c_empty:
                text = current["text"]
                current = r
                current["text"] = text + current["text"]
                is_c_empty = is_empty(current["text"])
            else:
                if current:
                    news_runs2.append(current)
                    current = r
                    is_c_empty = is_empty(current["text"])
                else:
                    current = r
                    is_c_empty = is_empty(current["text"])

        if current:
            if is_empty(current["text"]):
                if news_runs2:
                    news_runs2[-1]["text"] += current["text"]
                else:
                    news_runs2.append(current)
            else:
                news_runs2.append(current)

        new_runs = news_runs2

    return new_runs


def pptx_layout(filename):
    presentation = pptx.Presentation(filename)

    slide_height = presentation.slide_height
    slide_width = presentation.slide_width
    page_bbox = [slide_width, slide_height]
    page_bbox = norm_unit(page_bbox)
    layout = {"pages": [], "page_bbox": page_bbox}

    for i, slide in enumerate(presentation.slides):
        shape_infos = []
        for shape in _order_shapes(slide.shapes):
            # print('shape xml', etree.tostring(shape.element).decode())

            bbox = [shape.left, shape.top, shape.width, shape.height]
            shape_info = []
            if not shape.has_text_frame:
                continue

            # print('ori text', [shape.text_frame.text])
            whole_text = shape.text_frame.text.replace("\x0b", "\n")
            lines = whole_text.split("\n")
            lines_cnt = [len(line) for line in lines if len(line) > 0]

            for para in shape.text_frame.paragraphs:
                p_align = norm_align(para.alignment)
                p_line_spacing = para.line_spacing
                # p_space_after = para.space_after
                # p_space_before = para.space_before
                print("line_spacing", p_line_spacing, para.text)

                for run in para.runs:
                    if len(run.text) == 0:
                        continue

                    font_name = run.font.name
                    font_size = run.font.size
                    font_size = norm_unit(font_size)
                    color = run.font.color
                    color = str(color.rgb) if hasattr(color, "rgb") else None
                    color = norm_color(color)
                    info = {
                        "font_name": font_name,
                        "font_size": font_size,
                        "font_color": color,
                        "font_bold": run.font.bold,
                        "font_italic": run.font.italic,
                        "font_underline": run.font.underline,
                        "line_spacing": p_line_spacing,
                        # "space_after": p_space_after,
                        # "space_before": p_space_before,
                        "align": p_align,
                        "text": run.text,
                    }
                    shape_info.append(info)

            shape_texts = [info["text"] for info in shape_info]
            shapes_cnt = [len(info["text"]) for info in shape_info]
            assert sum(shapes_cnt) == sum(lines_cnt), (shapes_cnt, lines)

            if sum(shapes_cnt) == 0:
                continue

            lines_index = []
            n1 = len(shapes_cnt)
            n2 = len(lines_cnt)
            a1_i = 0
            a2_j = 0
            for _ in lines_cnt:
                found = False
                for ind in range(a1_i, n1):
                    if sum(shapes_cnt[a1_i : ind + 1]) == lines_cnt[a2_j]:
                        lines_index.append([a1_i, ind + 1])
                        a1_i = ind + 1
                        a2_j += 1
                        found = True
                        break

                assert found == True, "not found"

            shape_line_group = []
            for s, e in lines_index:
                new_runs = group_runs(shape_info[s:e], True)
                shape_line_group.append(new_runs)

            assert len(shape_line_group) == len(lines_cnt), "error in grouping"

            k = 0
            new_shape_line_group = []
            for line in lines:
                # empty line
                if len(line) == 0:
                    shape_info = [
                        {
                            "font_name": None,
                            "font_size": None,
                            "font_color": None,
                            "font_bold": None,
                            "font_italic": None,
                            "font_underline": None,
                            "line_spacing": None,
                            # "space_after": None,
                            # "space_before": None,
                            "align": None,
                            "text": "\n",
                        }
                    ]
                    new_shape_line_group.append(shape_info)
                else:
                    new_shape_line_group.append(shape_line_group[k])
                    k += 1

            bbox = norm_unit(bbox)
            shape_infos.append(
                {"bbox": bbox, "lines": new_shape_line_group, "whole_text": whole_text},
            )

        layout["pages"].append((i, shape_infos))

    return layout


def test1():
    filename = "./examples/docs/毛泽东课件.pptx"
    layout = pptx_layout(filename)
    with open("out.txt", "w") as fout:
        fout.write(json.dumps(layout, indent=2, ensure_ascii=False))


test1()
